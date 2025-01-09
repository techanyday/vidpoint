from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import logging

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def create_presentation(slides_content, output_dir="presentations", title="Video Summary"):
    """
    Create a PowerPoint presentation using python-pptx
    
    Args:
        slides_content (list): List of dictionaries containing slide content
        output_dir (str): Directory to save the presentation
        title (str): Title of the presentation
    
    Returns:
        str: Path to the created presentation file
    """
    try:
        # Create presentations directory if it doesn't exist
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            logger.info(f"Created output directory: {output_dir}")

        # Create presentation
        prs = Presentation()
        
        # Add slides
        for content in slides_content:
            layout = content.get('layout', 'CONTENT')
            
            if layout == 'TITLE':
                # Title slide
                slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
                title_shape = slide.shapes.title
                subtitle_shape = slide.placeholders[1]
                
                # Format title
                title_shape.text = content['title']
                title_para = title_shape.text_frame.paragraphs[0]
                title_para.font.size = Pt(48)  # Larger title
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.CENTER
                
                # Format subtitle
                subtitle_shape.text = content.get('content', '')
                sub_para = subtitle_shape.text_frame.paragraphs[0]
                sub_para.font.size = Pt(32)
                sub_para.alignment = PP_ALIGN.CENTER
                
            elif layout == 'BULLET_POINTS':
                # Content slide with single point
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
                
                # Format title
                if content.get('title'):
                    title_shape = slide.shapes.title
                    title_shape.text = content['title']
                    title_para = title_shape.text_frame.paragraphs[0]
                    title_para.font.size = Pt(40)
                    title_para.font.bold = True
                    title_para.alignment = PP_ALIGN.CENTER
                
                # Add point in larger font since it's just one
                if content.get('content'):
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear()
                    
                    # Center the text box vertically
                    body_shape.top = Inches(2.5)
                    
                    # Add the single point
                    for point in content['content']:
                        p = tf.add_paragraph()
                        p.text = str(point)
                        p.font.size = Pt(32)  # Larger font for single point
                        p.font.bold = None
                        p.alignment = PP_ALIGN.CENTER  # Center align
                        p.space_before = Pt(0)  # No extra space needed for single point
                        p.space_after = Pt(0)
                        p.level = 0
                        
            else:  # Default content slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                # Format title
                if content.get('title'):
                    title_shape = slide.shapes.title
                    title_shape.text = content['title']
                    title_para = title_shape.text_frame.paragraphs[0]
                    title_para.font.size = Pt(40)
                    title_para.font.bold = True
                    title_para.alignment = PP_ALIGN.CENTER
                
                # Format content
                if content.get('content'):
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear()
                    
                    for point in content['content']:
                        p = tf.add_paragraph()
                        p.text = str(point)
                        p.font.size = Pt(28)
                        p.alignment = PP_ALIGN.LEFT
                        p.level = 0
        
        # Save the presentation
        try:
            prs.save(os.path.join(output_dir, f"{title.replace(' ', '_')}.pptx"))
            logger.info(f"Successfully saved presentation to {output_dir}")
            return os.path.join(output_dir, f"{title.replace(' ', '_')}.pptx")
        except Exception as e:
            logger.error(f"Failed to save presentation: {str(e)}")
            return None
            
    except Exception as e:
        logger.error(f"Error creating presentation: {str(e)}")
        return None
