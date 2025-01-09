import os
from tkinter import Tk, filedialog
from pptx import Presentation


def create_presentation(summary):
    """
    Creates a PowerPoint presentation from the summary and allows the user to choose the save location.
    """
    # Create a new presentation
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Summary"
    content.text = summary

    # Open a folder selection dialog
    Tk().withdraw()  # Hide the main tkinter window
    folder_path = filedialog.askdirectory(title="Select Folder to Save Presentation")

    if folder_path:
        output_file = os.path.join(folder_path, "presentation.pptx")
        try:
            # Save the presentation in the chosen folder
            presentation.save(output_file)
            print(f"Presentation saved as {output_file}")
        except PermissionError:
            print(f"Permission denied: Unable to save to {output_file}. Make sure the file is not open.")
    else:
        print("No folder selected. Presentation not saved.")
